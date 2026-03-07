#!/usr/bin/env python3
import argparse
import json
import os
import sys


def main():
    parser = argparse.ArgumentParser(description="Build PPTX from autonomy research bundle.")
    parser.add_argument("--output", required=True, help="Output pptx path")
    parser.add_argument("--title", default="资料汇报", help="Deck title")
    parser.add_argument("--subtitle", default="", help="Deck subtitle")
    parser.add_argument("--slides-json", default="[]", help="Slides JSON array")
    args = parser.parse_args()

    try:
        from pptx import Presentation
    except Exception as exc:
        payload = {
            "ok": False,
            "reason": "python_pptx_missing",
            "error": str(exc)
        }
        print(json.dumps(payload, ensure_ascii=False))
        return 3

    try:
        slides = json.loads(args.slides_json or "[]")
    except Exception:
        slides = []
    if not isinstance(slides, list):
        slides = []

    output_path = os.path.abspath(args.output)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    presentation = Presentation()
    title_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_layout)
    if title_slide.shapes.title:
        title_slide.shapes.title.text = args.title or "资料汇报"
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = args.subtitle or ""

    content_layout = presentation.slide_layouts[1] if len(presentation.slide_layouts) > 1 else presentation.slide_layouts[0]
    for slide_data in slides[:12]:
        if not isinstance(slide_data, dict):
            continue
        title = str(slide_data.get("title", "")).strip() or "资料要点"
        bullets = slide_data.get("bullets", [])
        if not isinstance(bullets, list):
            bullets = []
        slide = presentation.slides.add_slide(content_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title
        body = slide.shapes.placeholders[1] if len(slide.shapes.placeholders) > 1 else None
        if body and hasattr(body, "text_frame"):
            frame = body.text_frame
            frame.clear()
            if bullets:
                frame.text = str(bullets[0])[:400]
                for bullet in bullets[1:8]:
                    paragraph = frame.add_paragraph()
                    paragraph.text = str(bullet)[:400]
                    paragraph.level = 0
            else:
                frame.text = "暂无补充内容"

    presentation.save(output_path)
    payload = {
        "ok": True,
        "reason": "pptx_created",
        "output": output_path,
        "slides": len(slides[:12]) + 1
    }
    print(json.dumps(payload, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
